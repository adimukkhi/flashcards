from flask import Flask, render_template, redirect, url_for, request
from flask_sqlalchemy import SQLAlchemy
from google import genai
from google.genai import types
from dotenv import load_dotenv
import os
import json
from datetime import datetime
import markdown
from pptx import Presentation

load_dotenv()
folder = "uploads"
app = Flask(__name__)

app.config["folder"] = folder
app.config["ASSIGN_FOLDER"] = folder
app.config["UPLOAD_FOLDER"] = folder
app.config["SQLALCHEMY_DATABASE_URI"] = "sqlite:///database.db"
app.config["SQLALCHEMY_TRACK_MODIFICATIONS"] = False

db = SQLAlchemy(app)
key = os.getenv("GEMINI_API_KEY")
client = genai.Client(api_key=key)

os.makedirs(folder, exist_ok=True)

class FlashcardDeck(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    title = db.Column(db.String(150), nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    cards = db.relationship('FlashcardItem', backref='deck', lazy=True, cascade="all, delete-orphan")

class FlashcardItem(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    deck_id = db.Column(db.Integer, db.ForeignKey('flashcard_deck.id'), nullable=False)
    question = db.Column(db.Text, nullable=False)
    answer = db.Column(db.Text, nullable=False)

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    slide_text_list = []
    for i, slide in enumerate(prs.slides):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())
        if slide_content:
            slide_text_list.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_content))
    return "\n\n".join(slide_text_list)

@app.template_filter('markdown')
def convert_markdown(text):
    return markdown.markdown(text, extensions=['fenced_code', 'tables'])

@app.route("/", methods=["GET", "POST"])
def index():
    prompt = ""
    number = 10
    active_deck_id = request.args.get("deck_id", type=int)
    all_decks = FlashcardDeck.query.order_by(FlashcardDeck.created_at.desc()).all()
    current_flashcards = []

    if active_deck_id:
        deck = db.session.get(FlashcardDeck, active_deck_id)
        if deck:
            current_flashcards = [{"question": c.question, "answer": c.answer} for c in deck.cards]

    if request.method == "POST":
        source = request.form.get("name", "").strip()
        number = request.form.get("number", "10")
        file = request.files.get("media_file")

        payload = []
        prompt = []

        if file and file.filename != "":
            path = os.path.join(app.config["folder"], file.filename)
            file.save(path)

            try:
                if file.filename.lower().endswith(('.ppt', '.pptx')):
                    pptx_text = extract_text_from_pptx(path)
                    payload.append(f"\n\nEXTRACTED SLIDE PRESENTATION CONTENT:\n{pptx_text}")
                else:
                    media = client.files.upload(file=path)
                    payload.append(media)
            finally:
                if os.path.exists(path):
                    os.remove(path)

        prompt.append(
            f"""You are an advanced educational processor for a study app.
            
            Analyze the attached source material (text, audio, or video) and execute two tasks:
            1. Generate a short, highly descriptive study title (3-5 words max) summarizing the core topic. Do not include quotes.
            2. Create exactly {number} comprehensive study flashcards testing the material. 

            Rules for cards: Clear specific questions, concise answers (1-2 sentences max), no concept repetitions.

            You must return your response using exactly this JSON structural dictionary frame:
            {{
                "title": "A Smart Descriptive Title Here",
                "flashcards": [
                    {{"question": "...", "answer": "..."}}
                ]
            }}"""
        )

        payload.extend(prompt)

        if source:
            payload.append(f"\n\nSOURCE TEXT:\n{source}")

    if request.method == "POST" and payload:
        response = client.models.generate_content(
            model="gemini-3.5-flash-lite",
            contents = payload,
            config = types.GenerateContentConfig(
                response_mime_type="application/json"
            )
        )

        try:
            raw = response.text
            payload = json.loads(raw)
            title = payload.get("title", "AI Deck")
            flashcards = payload.get("flashcards", [])
            new_deck = FlashcardDeck(title=title)
            db.session.add(new_deck)
            db.session.flush()

            for card in flashcards:
                item = FlashcardItem(
                    deck_id = new_deck.id,
                    question = card["question"],
                    answer = card["answer"]
                )
                db.session.add(item)

            db.session.commit()
            return redirect(url_for("index", deck_id=new_deck.id))
        except Exception as e:
            db.session.rollback()
            print(f"Error: {e}")
            current_flashcards = [{"question": "Oops, We ran into an error!", "answer": "Please try again :("}]

    return render_template("index.html", flashcards=current_flashcards, history_decks=all_decks, active_id=active_deck_id)

@app.route("/delete/<int:deck_id>", methods=["POST"])
def delete_deck(deck_id):
    target = db.session.get(FlashcardDeck, deck_id)
    if target:
        db.session.delete(target)
        db.session.commit()
    return redirect(url_for("index"))        

if __name__ == "__main__":
    app.run(debug=True)
